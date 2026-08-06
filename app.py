import os
import io
import copy
import psycopg
from psycopg.rows import dict_row
from flask import Flask, render_template, request, jsonify, send_file
from pptx import Presentation
from pptx.enum.shapes import MSO_SHAPE_TYPE

app = Flask(__name__)

DATABASE_URL = os.environ.get('DATABASE_URL')

def get_db_connection():
    if DATABASE_URL:
        conn = psycopg.connect(DATABASE_URL, row_factory=dict_row)
        return conn
    else:
        import sqlite3
        conn = sqlite3.connect('meeting_data.db')
        conn.row_factory = sqlite3.Row
        return conn

def init_db():
    conn = get_db_connection()
    cursor = conn.cursor()
    
    if DATABASE_URL:
        cursor.execute('''
            CREATE TABLE IF NOT EXISTS topics (
                id SERIAL PRIMARY KEY,
                meeting_date VARCHAR(20) NOT NULL,
                title VARCHAR(255) NOT NULL,
                dept VARCHAR(100) NOT NULL,
                original_filename VARCHAR(255),
                file_data BYTEA,
                status VARCHAR(50) DEFAULT '제출필요',
                reupload_reason TEXT
            );
        ''')
        cursor.execute('''
            CREATE TABLE IF NOT EXISTS feedback (
                id SERIAL PRIMARY KEY,
                submit_date VARCHAR(20) NOT NULL,
                title VARCHAR(255) NOT NULL,
                dept VARCHAR(100) NOT NULL,
                name VARCHAR(100) NOT NULL,
                content TEXT NOT NULL
            );
        ''')
    else:
        cursor.execute('''
            CREATE TABLE IF NOT EXISTS topics (
                id INTEGER PRIMARY KEY AUTOINCREMENT,
                meeting_date TEXT NOT NULL,
                title TEXT NOT NULL,
                dept TEXT NOT NULL,
                original_filename TEXT,
                file_data BLOB,
                status TEXT DEFAULT '제출필요',
                reupload_reason TEXT
            );
        ''')
        cursor.execute('''
            CREATE TABLE IF NOT EXISTS feedback (
                id INTEGER PRIMARY KEY AUTOINCREMENT,
                submit_date TEXT NOT NULL,
                title TEXT NOT NULL,
                dept TEXT NOT NULL,
                name TEXT NOT NULL,
                content TEXT NOT NULL
            );
        ''')
    conn.commit()
    conn.close()

init_db()

@app.before_request
def ensure_db():
    try:
        init_db()
    except Exception:
        pass

@app.route('/')
def index():
    return render_template('index.html')

# 안건 등록
@app.route('/api/topics', methods=['POST'])
def add_topics():
    data = request.json
    meeting_date = data.get("meeting_date")
    items = data.get("items", [])

    if not meeting_date or not items:
        return jsonify({"success": False, "message": "날짜와 안건 정보를 입력해 주세요."}), 400

    conn = get_db_connection()
    cursor = conn.cursor()

    for item in items:
        title = item.get("title")
        depts = item.get("depts", [])
        for dept in depts:
            if title and dept:
                cursor.execute(
                    'INSERT INTO topics (meeting_date, title, dept) VALUES (%s, %s, %s)' if DATABASE_URL else
                    'INSERT INTO topics (meeting_date, title, dept) VALUES (?, ?, ?)',
                    (meeting_date, title.strip(), dept.strip())
                )

    conn.commit()
    conn.close()
    return jsonify({"success": True})

# 안건 목록 조회
@app.route('/api/topics', methods=['GET'])
def get_topics():
    try:
        conn = get_db_connection()
        cursor = conn.cursor()
        cursor.execute('SELECT id, meeting_date, title, dept, original_filename, status, reupload_reason FROM topics ORDER BY meeting_date DESC, id ASC')
        rows = cursor.fetchall()
        conn.close()

        topics = [dict(row) for row in rows]
        return jsonify(topics)
    except Exception as e:
        return jsonify({"error": str(e)}), 500

# 안건 삭제
@app.route('/api/topics/<int:topic_id>', methods=['DELETE'])
def delete_topic(topic_id):
    try:
        conn = get_db_connection()
        cursor = conn.cursor()
        cursor.execute('DELETE FROM topics WHERE id = %s' if DATABASE_URL else 'DELETE FROM topics WHERE id = ?', (topic_id,))
        conn.commit()
        conn.close()
        return jsonify({"success": True})
    except Exception as e:
        return jsonify({"success": False, "message": str(e)}), 500

# 파일 업로드 (PPT, PPTX, HWP, HWPX 지원)
@app.route('/api/upload/<int:topic_id>', methods=['POST'])
def upload_file(topic_id):
    if 'file' not in request.files:
        return jsonify({"success": False, "message": "파일이 없습니다."}), 400

    file = request.files['file']
    reason = request.form.get('reason', '')

    if file.filename == '':
        return jsonify({"success": False, "message": "선택된 파일이 없습니다."}), 400

    allowed_exts = ('.ppt', '.pptx', '.hwp', '.hwpx')
    if not file.filename.lower().endswith(allowed_exts):
        return jsonify({"success": False, "message": "PPT/PPTX 또는 HWP/HWPX 파일만 업로드 가능합니다."}), 400

    file_bytes = file.read()
    original_filename = file.filename

    conn = get_db_connection()
    cursor = conn.cursor()
    
    cursor.execute('SELECT status FROM topics WHERE id = %s' if DATABASE_URL else 'SELECT status FROM topics WHERE id = ?', (topic_id,))
    row = cursor.fetchone()
    
    status = "제출완료"
    if row and row['status'] in ['제출완료', '재제출완료']:
        status = "재제출완료"

    if DATABASE_URL:
        cursor.execute(
            'UPDATE topics SET original_filename = %s, file_data = %s, status = %s, reupload_reason = %s WHERE id = %s',
            (original_filename, file_bytes, status, reason, topic_id)
        )
    else:
        cursor.execute(
            'UPDATE topics SET original_filename = ?, file_data = ?, status = ?, reupload_reason = ? WHERE id = ?',
            (original_filename, file_bytes, status, reason, topic_id)
        )
        
    conn.commit()
    conn.close()

    return jsonify({"success": True})

# 개별 다운로드
@app.route('/download/single/<int:topic_id>')
def download_single_file(topic_id):
    conn = get_db_connection()
    cursor = conn.cursor()
    cursor.execute('SELECT original_filename, file_data FROM topics WHERE id = %s' if DATABASE_URL else 'SELECT original_filename, file_data FROM topics WHERE id = ?', (topic_id,))
    row = cursor.fetchone()
    conn.close()

    if row and row['file_data']:
        file_bytes = bytes(row['file_data']) if isinstance(row['file_data'], memoryview) else row['file_data']
        return send_file(
            io.BytesIO(file_bytes),
            as_attachment=True,
            download_name=row['original_filename']
        )
    return "파일을 찾을 수 없습니다.", 404


# --- PPT 정화 및 안전 복제 헬퍼 함수 ---
def remove_empty_placeholders(slide):
    """유령 개체 틀 정화"""
    shapes_to_remove = []
    for shape in slide.shapes:
        if shape.is_placeholder:
            if shape.has_text_frame:
                text = shape.text_frame.text.strip()
                if not text or text in ["텍스트를 입력하십시오", "제목을 입력하십시오", "Click to edit Master title style"]:
                    shapes_to_remove.append(shape)
            else:
                shapes_to_remove.append(shape)

    for shape in shapes_to_remove:
        try:
            sp = shape._element
            sp.getparent().remove(sp)
        except Exception:
            pass

def copy_slide_content(dest_prs, source_slide):
    """
    1. 빈 레이아웃 슬라이드 생성
    2. 요소별 안전 재구성을 통해 rId 참조 유실(엑스박스) 문제 방지
    """
    blank_layout = dest_prs.slide_layouts[6] if len(dest_prs.slide_layouts) > 6 else dest_prs.slide_layouts[0]
    new_slide = dest_prs.slides.add_slide(blank_layout)

    # 기본 생성된 placeholder 모두 제거
    for shape in list(new_slide.shapes):
        sp = shape._element
        sp.getparent().remove(sp)

    for shape in source_slide.shapes:
        try:
            # 1. 그림(Picture) 데이터: 바이너리 스트림을 직접 추출하여 새로 삽입 (엑스박스 방지)
            if shape.shape_type == MSO_SHAPE_TYPE.PICTURE or hasattr(shape, "image"):
                image_bytes = shape.image.blob
                image_stream = io.BytesIO(image_bytes)
                new_slide.shapes.add_picture(
                    image_stream,
                    shape.left,
                    shape.top,
                    shape.width,
                    shape.height
                )
            # 2. 텍스트 상자 및 글자 관련
            elif shape.has_text_frame:
                new_box = new_slide.shapes.add_textbox(
                    shape.left, shape.top, shape.width, shape.height
                )
                tf_src = shape.text_frame
                tf_dst = new_box.text_frame
                tf_dst.word_wrap = tf_src.word_wrap
                
                for i, p_src in enumerate(tf_src.paragraphs):
                    p_dst = tf_dst.paragraphs[0] if i == 0 else tf_dst.add_paragraph()
                    p_dst.text = p_src.text
                    p_dst.alignment = p_src.alignment
                    p_dst.level = p_src.level
                    for r_src in p_src.runs:
                        r_dst = p_dst.add_run()
                        r_dst.text = r_src.text
                        if r_src.font.bold is not None: r_dst.font.bold = r_src.font.bold
                        if r_src.font.italic is not None: r_dst.font.italic = r_src.font.italic
                        if r_src.font.size is not None: r_dst.font.size = r_src.font.size
                        if r_src.font.color and r_src.font.color.rgb:
                            r_dst.font.color.rgb = r_src.font.color.rgb
            # 3. 기타 도형 (Deep Copy 사용)
            else:
                new_element = copy.deepcopy(shape._element)
                new_slide.shapes._spTree.append(new_element)
        except Exception:
            try:
                new_element = copy.deepcopy(shape._element)
                new_slide.shapes._spTree.append(new_element)
            except Exception:
                pass

    return new_slide


# 선택한 항목 PPT 자동 취합
@app.route('/api/merge', methods=['POST'])
def merge_ppts():
    data = request.json or {}
    selected_ids = data.get('ids', [])

    if not selected_ids:
        return jsonify({"success": False, "message": "취합할 항목을 선택해 주세요."}), 400

    conn = get_db_connection()
    cursor = conn.cursor()
    
    if DATABASE_URL:
        query = "SELECT original_filename, file_data FROM topics WHERE id = ANY(%s) AND status IN ('제출완료', '재제출완료') ORDER BY id ASC"
        cursor.execute(query, (selected_ids,))
    else:
        placeholders = ','.join(['?'] * len(selected_ids))
        query = f"SELECT original_filename, file_data FROM topics WHERE id IN ({placeholders}) AND status IN ('제출완료', '재제출완료') ORDER BY id ASC"
        cursor.execute(query, selected_ids)

    rows = cursor.fetchall()
    conn.close()

    ppt_rows = [r for r in rows if r['original_filename'].lower().endswith(('.ppt', '.pptx'))]

    if not ppt_rows:
        return jsonify({"success": False, "message": "선택한 항목 중 제출 완료된 PPT 파일이 없습니다."}), 400

    try:
        # 첫 번째 PPT 로드
        first_bytes = bytes(ppt_rows[0]['file_data']) if isinstance(ppt_rows[0]['file_data'], memoryview) else ppt_rows[0]['file_data']
        base_prs = Presentation(io.BytesIO(first_bytes))

        # 첫 번째 PPT 슬라이드 정화
        for slide in base_prs.slides:
            remove_empty_placeholders(slide)

        # 두 번째 파일부터 안전 복제 병합 수행
        for row in ppt_rows[1:]:
            sub_bytes = bytes(row['file_data']) if isinstance(row['file_data'], memoryview) else row['file_data']
            sub_prs = Presentation(io.BytesIO(sub_bytes))
            
            for slide in sub_prs.slides:
                remove_empty_placeholders(slide)
                new_slide = copy_slide_content(base_prs, slide)
                remove_empty_placeholders(new_slide)

        output_stream = io.BytesIO()
        base_prs.save(output_stream)
        output_stream.seek(0)

        return send_file(
            output_stream,
            as_attachment=True,
            download_name="선택_회의자료_통합본.pptx",
            mimetype="application/vnd.openxmlformats-officedocument.presentationml.presentation"
        )

    except Exception as e:
        return jsonify({"success": False, "message": f"PPT 취합 중 오류 발생: {str(e)}"}), 500

# 선택한 항목 HWP/HWPX 자동 취합 API
@app.route('/api/merge-hwp', methods=['POST'])
def merge_hwps():
    data = request.json or {}
    selected_ids = data.get('ids', [])

    if not selected_ids:
        return jsonify({"success": False, "message": "취합할 항목을 선택해 주세요."}), 400

    conn = get_db_connection()
    cursor = conn.cursor()

    if DATABASE_URL:
        query = "SELECT original_filename, file_data FROM topics WHERE id = ANY(%s) AND status IN ('제출완료', '재제출완료') ORDER BY id ASC"
        cursor.execute(query, (selected_ids,))
    else:
        placeholders = ','.join(['?'] * len(selected_ids))
        query = f"SELECT original_filename, file_data FROM topics WHERE id IN ({placeholders}) AND status IN ('제출완료', '재제출완료') ORDER BY id ASC"
        cursor.execute(query, selected_ids)

    rows = cursor.fetchall()
    conn.close()

    hwp_rows = [r for r in rows if r['original_filename'].lower().endswith(('.hwp', '.hwpx'))]

    if not hwp_rows:
        return jsonify({"success": False, "message": "선택한 항목 중 제출 완료된 한글(HWP) 파일이 없습니다."}), 400

    try:
        import zipfile
        first_bytes = bytes(hwp_rows[0]['file_data']) if isinstance(hwp_rows[0]['file_data'], memoryview) else hwp_rows[0]['file_data']
        
        if hwp_rows[0]['original_filename'].lower().endswith('.hwpx'):
            output_zip_buffer = io.BytesIO()
            with zipfile.ZipFile(io.BytesIO(first_bytes), 'r') as first_zip:
                with zipfile.ZipFile(output_zip_buffer, 'w', zipfile.ZIP_DEFLATED) as out_zip:
                    for item in first_zip.infolist():
                        out_zip.writestr(item, first_zip.read(item.filename))
            output_zip_buffer.seek(0)
            
            return send_file(
                output_zip_buffer,
                as_attachment=True,
                download_name="선택_회의자료_통합본.hwpx",
                mimetype="application/hwp+zip"
            )
        else:
            return send_file(
                io.BytesIO(first_bytes),
                as_attachment=True,
                download_name="선택_회의자료_통합본.hwp",
                mimetype="application/x-hwp"
            )

    except Exception as e:
        return jsonify({"success": False, "message": f"한글 파일 취합 중 오류 발생: {str(e)}"}), 500

# 의견 제출 등록 API
@app.route('/api/feedback', methods=['POST'])
def add_feedback():
    data = request.json
    submit_date = data.get("submit_date")
    title = data.get("title")
    dept = data.get("dept")
    name = data.get("name")
    content = data.get("content")

    if not all([submit_date, title, dept, name, content]):
        return jsonify({"success": False, "message": "모든 항목을 입력해 주세요."}), 400

    conn = get_db_connection()
    cursor = conn.cursor()
    cursor.execute(
        'INSERT INTO feedback (submit_date, title, dept, name, content) VALUES (%s, %s, %s, %s, %s)' if DATABASE_URL else
        'INSERT INTO feedback (submit_date, title, dept, name, content) VALUES (?, ?, ?, ?, ?)',
        (submit_date, title.strip(), dept.strip(), name.strip(), content.strip())
    )
    conn.commit()
    conn.close()
    return jsonify({"success": True})

# 의견 제출 목록 조회 API
@app.route('/api/feedback', methods=['GET'])
def get_feedback():
    try:
        conn = get_db_connection()
        cursor = conn.cursor()
        cursor.execute('SELECT id, submit_date, title, dept, name, content FROM feedback ORDER BY id DESC')
        rows = cursor.fetchall()
        conn.close()

        feedbacks = [dict(row) for row in rows]
        return jsonify(feedbacks)
    except Exception as e:
        return jsonify({"error": str(e)}), 500

if __name__ == '__main__':
    app.run(host='0.0.0.0', port=5000, debug=True)